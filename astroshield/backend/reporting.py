"""Presentation export helpers for AstroShield briefings."""
from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Any, Dict, Iterable

from pptx import Presentation
from pptx.util import Inches, Pt


def build_simulation_briefing(
    simulation: Dict[str, Any],
    *,
    generated_at: datetime | None = None,
    author: str | None = None,
) -> bytes:
    """Create a lightweight mission briefing deck for a simulation payload."""

    generated_at = generated_at or datetime.utcnow()

    prs = Presentation()
    _populate_title_slide(prs, simulation, generated_at, author)
    _populate_summary_slide(prs, simulation)
    _populate_impact_slide(prs, simulation)
    _populate_environment_slide(prs, simulation)

    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream.read()


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _populate_title_slide(prs: Presentation, simulation: Dict[str, Any], generated_at: datetime, author: str | None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    neo = simulation.get("neo_reference", {})
    inputs = simulation.get("inputs", {})

    friendly_id = neo.get("friendly_id") or inputs.get("asteroid_id") or "Asteroid"
    name = neo.get("name") or friendly_id
    title.text = f"AstroShield Mission Briefing: {name}"

    lines = [
        f"Generated {generated_at.strftime('%Y-%m-%d %H:%M UTC')}",
        f"Asteroid ID: {inputs.get('asteroid_id', friendly_id)}",
    ]
    if author:
        lines.append(f"Prepared for {author}")
    subtitle.text = "\n".join(lines)


def _populate_summary_slide(prs: Presentation, simulation: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Snapshot"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    neo = simulation.get("neo_reference", {})
    energy = simulation.get("energy", {})
    effects = simulation.get("impact_effects", {})

    bullets = [
        f"Source: {neo.get('source', 'mock').title()}",
        f"Diameter: {_format_range(neo.get('diameter_range_m'), units='m', default=_format_number(neo.get('diameter_m'), 'm'))}",
        f"Velocity: {_format_number(energy.get('effective_velocity_ms') / 1000 if energy.get('effective_velocity_ms') else neo.get('velocity_kms'), 'km/s')}",
        f"Impact energy: {_format_number(energy.get('energy_mt'), 'Mt TNT')}",
        f"Crater size: {_format_number(effects.get('crater_diameter_km'), 'km')}",
        f"Seismic magnitude: {_format_number(effects.get('seismic_magnitude'))}",
    ]

    for text in bullets:
        if text is None:
            continue
        p = body.add_paragraph()
        p.text = text
        p.level = 0


def _populate_impact_slide(prs: Presentation, simulation: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title + Content
    slide.shapes.title.text = "Impact Metrics"

    rows = [
        ("Mass", _format_number(simulation.get("energy", {}).get("mass_kg"), "kg")),
        ("Kinetic Energy", _format_number(simulation.get("energy", {}).get("energy_joules"), "J")),
        ("Crater Diameter", _format_number(simulation.get("impact_effects", {}).get("crater_diameter_km"), "km")),
        ("Seismic Magnitude", _format_number(simulation.get("impact_effects", {}).get("seismic_magnitude"))),
        ("Tsunami Risk", _format_boolean(simulation.get("environment", {}).get("tsunami_risk"))),
    ]

    columns = 2
    table = slide.shapes.add_table(len(rows) + 1, columns, Inches(0.5), Inches(1.5), Inches(9.0), Inches(4.0)).table
    table.columns[0].width = Inches(4.0)
    table.columns[1].width = Inches(5.0)

    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"

    for idx, (label, value) in enumerate(rows, start=1):
        table.cell(idx, 0).text = label
        table.cell(idx, 1).text = value

    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(18)


def _populate_environment_slide(prs: Presentation, simulation: Dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Environment & Orbit"
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    inputs = simulation.get("inputs", {})
    environment = simulation.get("environment", {})
    orbital = simulation.get("orbital_solution", {})

    bullets = [
        f"Impact coordinates: {inputs.get('impact_lat', 0):.2f}°, {inputs.get('impact_lon', 0):.2f}°",
        f"Elevation: {_format_number(environment.get('elevation_m'), 'm')}",
        f"Coastal zone: {_format_boolean(environment.get('is_coastal_zone'))}",
        f"Seismic risk: {environment.get('seismic_zone_risk', 'Unknown')}",
        f"Baseline MOID: {_format_number(orbital.get('baseline_moid_km'), 'km')}",
        f"Deflected MOID: {_format_number(orbital.get('deflected_moid_km'), 'km')}",
        f"MOID change: {_format_number(orbital.get('moid_change_km'), 'km')}",
    ]

    for text in bullets:
        p = body.add_paragraph()
        p.text = text
        p.level = 0

    if environment.get("tectonic_summary"):
        p = body.add_paragraph()
        p.text = f"Tectonic summary: {environment['tectonic_summary']}"
        p.level = 0


# ---------------------------------------------------------------------------
# Formatting helpers
# ---------------------------------------------------------------------------

def _format_number(value: Any, units: str | None = None, precision: int = 2, default: str = "—") -> str:
    try:
        if value is None:
            raise ValueError
        number = float(value)
        if abs(number) >= 1_000_000_000:
            formatted = f"{number/1_000_000_000:.{precision}f}B"
        elif abs(number) >= 1_000_000:
            formatted = f"{number/1_000_000:.{precision}f}M"
        elif abs(number) >= 1_000:
            formatted = f"{number/1_000:.{precision}f}k"
        else:
            formatted = f"{number:.{precision}f}"
        return f"{formatted}{(' ' + units) if units else ''}"
    except (TypeError, ValueError):
        return default


def _format_range(range_tuple: Any, *, units: str, default: str = "—") -> str:
    if not isinstance(range_tuple, Iterable):
        return default
    try:
        lower, upper = range_tuple
        lower_text = _format_number(lower, units, default="—") if lower is not None else "—"
        upper_text = _format_number(upper, units, default="—") if upper is not None else "—"
        if lower_text == upper_text:
            return lower_text
        return f"{lower_text} – {upper_text}"
    except (ValueError, TypeError):
        return default


def _format_boolean(value: Any) -> str:
    if value is None:
        return "Unknown"
    return "Yes" if bool(value) else "No"
